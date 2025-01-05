#from fpdf import FPDF
from PyPDF2 import PdfReader
from pptx import Presentation

def extract_ppt_text(file):
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text

def extract_pdf_text(file):
    pdf_reader = PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def split_text(text, chunk_size=500):
    words = text.split()
    chunks = [" ".join(words[i:i+chunk_size]) for i in range(0, len(words), chunk_size)]
    return chunks



